from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

_TEMPLATE = Path(__file__).parent / "templates" / "slides_template.pptx"

_IMG_LEFT  = Inches(5.2)
_IMG_TOP   = Inches(1.5)
_IMG_WIDTH = Inches(4.3)


def create_presentation(title: str, slides: str, filename: str, output_path: str | None = None) -> str:
  
    base = Path(output_path) if output_path else Path("output")
    base.mkdir(parents=True, exist_ok=True)
    output = base / f"{filename}.pptx"

    prs = Presentation(_TEMPLATE)

    # 1. Modificar la portada (slide 0)
    title_slide = prs.slides[0]
    title_slide.shapes.title.text = title
    # Limpiar subtítulo si existe
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = ""

    # 2. Obtener el diseño (slide layout) de la diapositiva de contenido (slide 1)
    if len(prs.slides) < 2:
        raise ValueError("La plantilla debe tener al menos dos diapositivas: portada y contenido.")
    
    content_slide = prs.slides[1]
    content_layout = content_slide.slide_layout  # diseño subyacente

    # 3. Eliminar la diapositiva de contenido original (para que no aparezca en el final)
    slides_list = prs.slides._sldIdLst
    slide_to_remove = slides_list[1]  # índice 1 = segunda diapositiva
    rId = slide_to_remove.get(qn("r:id"))
    prs.part.drop_rel(rId)
    slides_list.remove(slide_to_remove)

    # 4. Multiplicar: crear nuevas diapositivas de contenido según los elementos de 'slides'
    for item in slides.split("|"):
        if not item.strip():
            continue
        parts = item.split(":")
        slide = prs.slides.add_slide(content_layout)

        # Título de sección (primer placeholder de título)
        if slide.shapes.title:
            slide.shapes.title.text = parts[0].strip()

        # Contenido de texto (normalmente el segundo placeholder)
        if len(parts) > 1:
            texto = parts[1].strip()
            # Buscar el placeholder de contenido (puede ser el índice 1 o un shape de tipo BODY)
            try:
                slide.placeholders[1].text = texto
            except IndexError:
                # Alternativa: encontrar el primer placeholder que no sea título
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 2:  # tipo BODY
                        shape.text = texto
                        break

        # Imagen opcional (tercera parte)
        if len(parts) > 2:
            img_path = Path(parts[2].strip())
            if img_path.is_file():
                slide.shapes.add_picture(str(img_path), _IMG_LEFT, _IMG_TOP, width=_IMG_WIDTH)

    prs.save(output)
    return str(output.resolve())