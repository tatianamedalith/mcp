from pptx import Presentation

def create_presentation(title: str, slides: str, filename: str) -> str:
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generado por Lexi"

    for item in slides.split("|"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        partes = item.split(":")
        slide.shapes.title.text = partes[0].strip()
        if len(partes) > 1:
            slide.placeholders[1].text = partes[1].strip()

    output = f"{filename}.pptx"
    prs.save(output)
    return f"Presentación '{output}' creada."